import base64
import os
import re
import zipfile
import xml.etree.ElementTree as ET
from html import escape

from .page_preview_text import build_html_pages

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None


def extract_pptx_text(file_path: str, max_chars: int = 9000) -> str:
    texts: list[str] = []
    try:
        with zipfile.ZipFile(file_path) as zf:
            slide_names = sorted(
                [name for name in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml", name)]
            )
            for slide in slide_names:
                with zf.open(slide) as file_obj:
                    root = ET.fromstring(file_obj.read())
                for node in root.iter():
                    if node.tag.endswith("}t") and node.text:
                        texts.append(node.text)
                if sum(len(text) for text in texts) >= max_chars:
                    break
    except Exception:
        return ""
    return " ".join(texts)[:max_chars]


class PresentationPreviewService:
    EMU_PER_PIXEL = 9525.0

    def __init__(self, controller):
        self._controller = controller

    def get_preview_src(self, file_id: str, file_path: str, page: int) -> str:
        if not file_id or not file_path:
            return ""

        cached = self._controller._non_pdf_preview_cache.get(file_id)
        if cached is None:
            cached = self._build_preview_pages(file_path)
            if cached:
                self._controller._non_pdf_preview_cache[file_id] = cached
                self._controller._total_pages_cache[file_id] = max(1, len(cached))

        if not cached:
            return ""

        page_idx = max(1, int(page or 1)) - 1
        page_idx = min(page_idx, max(0, len(cached) - 1))
        return cached[page_idx]

    def extract_slide_text(
        self, file_path: str, page: int, max_chars: int = 7000
    ) -> str:
        if not file_path or not os.path.isfile(file_path) or Presentation is None:
            return ""

        try:
            presentation = Presentation(file_path)
        except Exception:
            return ""

        slides = list(presentation.slides)
        if not slides:
            return ""
        page_idx = max(0, min(len(slides) - 1, int(page or 1) - 1))
        texts: list[str] = []
        for shape in slides[page_idx].shapes:
            self._collect_shape_text(shape, texts)
            if sum(len(item) for item in texts) >= max_chars:
                break
        return " ".join(part for part in texts if part).strip()[:max_chars]

    def _build_preview_pages(self, file_path: str) -> list[str]:
        if not file_path or not os.path.isfile(file_path) or Presentation is None:
            return []

        try:
            presentation = Presentation(file_path)
        except Exception:
            return []

        slide_width = max(1, int(getattr(presentation, "slide_width", 1) or 1))
        slide_height = max(1, int(getattr(presentation, "slide_height", 1) or 1))
        page_contents: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            content = self._render_slide(slide, slide_index, slide_width, slide_height)
            page_contents.append(content)
        return build_html_pages(
            page_contents,
            body_class="pptx-preview-body",
            inline_script=self._build_preview_script(),
        )

    def _render_slide(self, slide, slide_index: int, slide_width: int, slide_height: int) -> str:
        slide_width_px = self._emu_to_px(slide_width)
        slide_height_px = self._emu_to_px(slide_height)
        background_style = self._get_slide_background_style(slide)
        elements: list[str] = []
        for z_index, shape in enumerate(slide.shapes, start=1):
            rendered = self._render_shape(shape, slide_width, slide_height, z_index)
            if rendered:
                elements.append(rendered)

        if not elements:
            elements.append(
                "<div class='pptx-preview-empty'>This slide has no previewable content.</div>"
            )

        return (
            "<section class='pptx-preview-root'>"
            "<div class='pptx-preview-shell'>"
            "<div class='pptx-preview-canvas-wrap'>"
            "<div class='pptx-preview-canvas-scale' data-pptx-scale='1'>"
            f"<div class='pptx-preview-stage' data-stage-width='{slide_width_px}' data-stage-height='{slide_height_px}' style='width:{slide_width_px:.2f}px;height:{slide_height_px:.2f}px;{background_style}'>"
            f"<div class='pptx-preview-slide-number'>Slide {slide_index} · Ctrl + Wheel Zoom</div>"
            f"{''.join(elements)}"
            "</div>"
            "</div>"
            "</div>"
            "</div>"
            "</section>"
        )

    def _render_shape(self, shape, slide_width: int, slide_height: int, z_index: int) -> str:
        try:
            shape_type = getattr(shape, "shape_type", None)
            if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.GROUP:
                return self._render_group_shape(shape, slide_width, slide_height, z_index)
        except Exception:
            pass

        if getattr(shape, "has_table", False):
            return self._render_table_shape(shape, slide_width, slide_height, z_index)

        if hasattr(shape, "image"):
            return self._render_picture_shape(shape, slide_width, slide_height, z_index)

        if getattr(shape, "has_text_frame", False):
            return self._render_text_shape(shape, slide_width, slide_height, z_index)

        box_style = self._build_shape_box_style(shape, slide_width, slide_height, z_index)
        if not box_style:
            return ""
        fill_style = self._extract_fill_style(getattr(shape, "fill", None), "background")
        line_style = self._extract_line_style(getattr(shape, "line", None))
        if not fill_style and not line_style:
            return ""
        return f"<div class='pptx-preview-shape' style='{box_style}{fill_style}{line_style}'></div>"

    def _render_group_shape(self, shape, slide_width: int, slide_height: int, z_index: int) -> str:
        group_text: list[str] = []
        for item in getattr(shape, "shapes", []):
            if getattr(item, "has_text_frame", False):
                text = self._extract_text_frame_html(item.text_frame)
                if text:
                    group_text.append(text)

        if not group_text:
            return ""

        box_style = self._build_shape_box_style(shape, slide_width, slide_height, z_index)
        fill_style = self._extract_fill_style(getattr(shape, "fill", None), "background")
        line_style = self._extract_line_style(getattr(shape, "line", None))
        return (
            f"<div class='pptx-preview-element pptx-preview-text' style='{box_style}{fill_style}{line_style}'>"
            f"{''.join(group_text)}"
            "</div>"
        )

    def _render_text_shape(self, shape, slide_width: int, slide_height: int, z_index: int) -> str:
        content = self._extract_text_frame_html(shape.text_frame)
        if not content:
            return ""
        box_style = self._build_shape_box_style(shape, slide_width, slide_height, z_index)
        text_frame_style = self._build_text_frame_style(shape.text_frame)
        fill_style = self._extract_fill_style(getattr(shape, "fill", None), "background")
        line_style = self._extract_line_style(getattr(shape, "line", None))
        return (
            f"<div class='pptx-preview-element pptx-preview-text' style='{box_style}{text_frame_style}{fill_style}{line_style}'>"
            f"{content}"
            "</div>"
        )

    def _render_picture_shape(self, shape, slide_width: int, slide_height: int, z_index: int) -> str:
        try:
            image = shape.image
            mime_type = getattr(image, "content_type", "") or "image/png"
            encoded = base64.b64encode(image.blob).decode("ascii")
        except Exception:
            return ""

        box_style = self._build_shape_box_style(shape, slide_width, slide_height, z_index)
        return (
            f"<div class='pptx-preview-element pptx-preview-picture' style='{box_style}'>"
            f"<img alt='' src='data:{mime_type};base64,{encoded}' />"
            "</div>"
        )

    def _render_table_shape(self, shape, slide_width: int, slide_height: int, z_index: int) -> str:
        try:
            table = shape.table
        except Exception:
            return ""

        rows_html: list[str] = []
        for row in table.rows:
            cells_html: list[str] = []
            for cell in row.cells:
                cell_fill = self._extract_fill_style(getattr(cell, "fill", None), "background")
                cell_text = escape((getattr(cell, "text", "") or "").strip()).replace("\n", "<br/>")
                if not cell_text:
                    cell_text = "&nbsp;"
                cells_html.append(
                    f"<td style='{cell_fill}'>{cell_text}</td>"
                )
            rows_html.append(f"<tr>{''.join(cells_html)}</tr>")

        if not rows_html:
            return ""

        box_style = self._build_shape_box_style(shape, slide_width, slide_height, z_index)
        line_style = self._extract_line_style(getattr(shape, "line", None))
        return (
            f"<div class='pptx-preview-element pptx-preview-table-wrap' style='{box_style}{line_style}'>"
            f"<table class='pptx-preview-table'>{''.join(rows_html)}</table>"
            "</div>"
        )

    def _build_text_frame_style(self, text_frame) -> str:
        if text_frame is None:
            return ""
        padding_top = self._emu_to_px(getattr(text_frame, "margin_top", 0) or 0)
        padding_right = self._emu_to_px(getattr(text_frame, "margin_right", 0) or 0)
        padding_bottom = self._emu_to_px(getattr(text_frame, "margin_bottom", 0) or 0)
        padding_left = self._emu_to_px(getattr(text_frame, "margin_left", 0) or 0)
        return (
            f"padding:{padding_top:.2f}px {padding_right:.2f}px {padding_bottom:.2f}px {padding_left:.2f}px;"
        )

    def _extract_text_frame_html(self, text_frame) -> str:
        paragraphs_html: list[str] = []
        for paragraph in getattr(text_frame, "paragraphs", []):
            paragraph_html = self._render_paragraph(paragraph)
            if paragraph_html:
                paragraphs_html.append(paragraph_html)
        return "".join(paragraphs_html)

    def _collect_shape_text(self, shape, texts: list[str]) -> None:
        if getattr(shape, "has_text_frame", False):
            for paragraph in getattr(shape.text_frame, "paragraphs", []):
                paragraph_text = (getattr(paragraph, "text", "") or "").strip()
                if paragraph_text:
                    texts.append(" ".join(paragraph_text.split()))

        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = (getattr(cell, "text", "") or "").strip()
                        if cell_text:
                            texts.append(" ".join(cell_text.split()))
            except Exception:
                pass

        try:
            shape_type = getattr(shape, "shape_type", None)
            if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.GROUP:
                for item in getattr(shape, "shapes", []):
                    self._collect_shape_text(item, texts)
        except Exception:
            pass

    def _render_paragraph(self, paragraph) -> str:
        runs_html: list[str] = []
        for run in getattr(paragraph, "runs", []):
            rendered = self._render_run(run)
            if rendered:
                runs_html.append(rendered)

        if not runs_html:
            plain_text = escape((getattr(paragraph, "text", "") or "").strip()).replace("\n", "<br/>")
            if plain_text:
                runs_html.append(plain_text)

        if not runs_html:
            return ""

        align_style = self._get_paragraph_alignment_style(paragraph)
        level = max(0, int(getattr(paragraph, "level", 0) or 0))
        indent_style = f"padding-left:{0.8 + level * 1.1}rem;" if level else ""
        return f"<p style='{align_style}{indent_style}'>{''.join(runs_html)}</p>"

    def _render_run(self, run) -> str:
        text = escape(getattr(run, "text", "") or "").replace("\n", "<br/>")
        if not text:
            return ""

        font = getattr(run, "font", None)
        style_tokens: list[str] = []
        if font is not None:
            if getattr(font, "name", None):
                style_tokens.append(f"font-family:{escape(str(font.name))},sans-serif;")
            size = getattr(font, "size", None)
            if size is not None:
                try:
                    px_size = max(10.0, min(72.0, float(size.pt) * 96.0 / 72.0))
                    style_tokens.append(f"font-size:{px_size:.2f}px;")
                except Exception:
                    pass
            color = self._extract_color_value(getattr(font, "color", None))
            if color:
                style_tokens.append(f"color:{color};")
            if getattr(font, "bold", False):
                style_tokens.append("font-weight:700;")
            if getattr(font, "italic", False):
                style_tokens.append("font-style:italic;")
            if getattr(font, "underline", False):
                style_tokens.append("text-decoration:underline;")

        hyperlink = ""
        try:
            hyperlink = (getattr(getattr(run, "hyperlink", None), "address", "") or "").strip()
        except Exception:
            hyperlink = ""

        style_attr = "".join(style_tokens)
        if hyperlink:
            return (
                f"<a class='pptx-preview-link' href='{escape(hyperlink)}' target='_blank' rel='noopener noreferrer' style='{style_attr}'>"
                f"{text}</a>"
            )
        return f"<span style='{style_attr}'>{text}</span>"

    def _get_slide_background_style(self, slide) -> str:
        fill = getattr(getattr(slide, "background", None), "fill", None)
        color = self._extract_fill_color_value(fill)
        if color:
            return f"background:{color};"
        return "background:linear-gradient(180deg,#ffffff 0%,#f4f7fb 100%);"

    def _build_shape_box_style(
        self,
        shape,
        slide_width: int,
        slide_height: int,
        z_index: int,
    ) -> str:
        try:
            left = self._emu_to_px(getattr(shape, "left", 0))
            top = self._emu_to_px(getattr(shape, "top", 0))
            width = self._emu_to_px(max(1, int(getattr(shape, "width", 1) or 1)))
            height = self._emu_to_px(max(1, int(getattr(shape, "height", 1) or 1)))
        except Exception:
            return ""

        style = (
            f"left:{left:.2f}px;top:{top:.2f}px;width:{width:.2f}px;height:{height:.2f}px;z-index:{z_index};"
        )
        rotation = getattr(shape, "rotation", 0) or 0
        try:
            rotation_value = float(rotation)
        except Exception:
            rotation_value = 0.0
        if abs(rotation_value) > 0.01:
            style += f"transform:rotate({rotation_value:.2f}deg);transform-origin:center center;"
        return style

    def _emu_to_px(self, value: int | float) -> float:
        return max(0.0, float(value or 0.0) / self.EMU_PER_PIXEL)

    def _extract_fill_style(self, fill, property_name: str) -> str:
        color = self._extract_fill_color_value(fill)
        if not color:
            return ""
        return f"{property_name}:{color};"

    def _extract_line_style(self, line) -> str:
        try:
            color_format = getattr(line, "color", None)
        except Exception:
            color_format = None
        color = self._extract_color_value(color_format)
        if not color:
            return ""
        width = "1px"
        try:
            line_width = getattr(line, "width", None)
            if line_width is not None:
                width = f"{max(1.0, min(6.0, float(line_width.pt))):.2f}px"
        except Exception:
            width = "1px"
        return f"border:{width} solid {color};"

    @staticmethod
    def _extract_color_value(color_format) -> str:
        if color_format is None:
            return ""
        try:
            rgb = getattr(color_format, "rgb", None)
        except Exception:
            rgb = None
        if rgb:
            return f"#{rgb}"
        return ""

    def _extract_fill_color_value(self, fill) -> str:
        if fill is None:
            return ""
        try:
            color_format = getattr(fill, "fore_color", None)
        except Exception:
            color_format = None
        return self._extract_color_value(color_format)

    def _build_preview_script(self) -> str:
        return """
(function () {
    const shell = document.querySelector('.pptx-preview-shell');
    const canvas = document.querySelector('.pptx-preview-canvas-scale');
    const stage = document.querySelector('.pptx-preview-stage');
    if (!shell || !canvas || !stage) {
        return;
    }

    let zoom = 1;

    function clamp(value, min, max) {
        return Math.max(min, Math.min(max, value));
    }

    function applyZoom(nextZoom) {
        zoom = clamp(nextZoom, 0.35, 3);
        canvas.style.zoom = String(zoom);
        canvas.setAttribute('data-pptx-scale', zoom.toFixed(2));
    }

    function fitIfNeeded() {
        const stageWidth = parseFloat(stage.getAttribute('data-stage-width') || '0');
        const viewportWidth = shell.clientWidth - 24;
        if (!stageWidth || !viewportWidth) {
            applyZoom(1);
            return;
        }
        if (stageWidth <= viewportWidth) {
            applyZoom(1);
            return;
        }
        applyZoom(viewportWidth / stageWidth);
    }

    shell.addEventListener('wheel', function (event) {
        if (!event.ctrlKey) {
            return;
        }
        event.preventDefault();
        const factor = event.deltaY < 0 ? 1.08 : 0.92;
        applyZoom(zoom * factor);
    }, { passive: false });

    window.addEventListener('resize', fitIfNeeded);
    fitIfNeeded();
})();
        """.strip()

    @staticmethod
    def _get_paragraph_alignment_style(paragraph) -> str:
        alignment = getattr(paragraph, "alignment", None)
        if alignment is None:
            return ""
        label = str(alignment).upper()
        if "CENTER" in label:
            return "text-align:center;"
        if "RIGHT" in label:
            return "text-align:right;"
        if "JUSTIFY" in label:
            return "text-align:justify;"
        return "text-align:left;"